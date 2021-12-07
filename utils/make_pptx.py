from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import glob

TARGET_DIR = '/Users/y1u0d2/desktop/Lab/result/nnp-train/two-three-body/03/01/error/r2'
energy_pic = glob.glob(f'{TARGET_DIR}/energy/*.png')
force_pic = glob.glob(f'{TARGET_DIR}/force/*.png')
energy_pic.sort()
force_pic.sort()

IMG_DISPLAY_HEIGHT = Inches(3)  # 画像サイズ
PIC_PER_PAGE = 4
OUTPUT_FILE_NAME = f"/Users/y1u0d2/desktop/Lab/result/nnp-train/two-three-body/03/01/r2.pptx"  # 出力ファイル名

slides=0
if len(energy_pic) % 2 == 0:
    slides = len(energy_pic)/2
else:
    slides = (len(energy_pic)+1)/2

prs = Presentation()
SLIDE_WIDTH = prs.slide_width
SLIDE_HEIGHT = prs.slide_height

for i in range(int(slides)):
    # slide_layouts:6 blank
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    energy_pic_1 = energy_pic[i*2]
    force_pic_1 = force_pic[i*2]

    if not i*2+1 >= len(energy_pic):
        energy_pic_2 = energy_pic[i * 2 + 1]
        force_pic_2 = force_pic[i*2+1]

    im = Image.open(energy_pic_1)
    im_width, im_height = im.size
    aspect_ratio = im_width / im_height

    img_display_height = IMG_DISPLAY_HEIGHT
    img_display_width = aspect_ratio * img_display_height

    # left top
    left = (SLIDE_WIDTH / 2 - img_display_width)
    top = (SLIDE_HEIGHT / 2 - img_display_height)
    slide.shapes.add_picture(energy_pic_1, left, top, height=IMG_DISPLAY_HEIGHT)
    # left bottom
    if not i * 2 + 1 >= len(energy_pic):
        left = (SLIDE_WIDTH / 2 - img_display_width)
        top = (SLIDE_HEIGHT / 2)
        slide.shapes.add_picture(energy_pic_2, left, top, height=IMG_DISPLAY_HEIGHT)
    # right top
    left = (SLIDE_WIDTH / 2)
    top = (SLIDE_HEIGHT / 2 - img_display_height)
    slide.shapes.add_picture(force_pic_1, left, top, height=IMG_DISPLAY_HEIGHT)
    # right bottom
    if not i * 2 + 1 >= len(energy_pic):
        left = (SLIDE_WIDTH / 2)
        top = (SLIDE_HEIGHT / 2)
        slide.shapes.add_picture(force_pic_2, left, top, height=IMG_DISPLAY_HEIGHT)

prs.save(OUTPUT_FILE_NAME)
# TODO convert pdf