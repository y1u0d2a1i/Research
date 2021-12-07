from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import glob

TARGET_DIR = '/Users/y1u0d2/desktop/Lab/result/nnp-train/20211130/pic/r2'
energy_pic = glob.glob(f'{TARGET_DIR}/energy/*.png')
force_pic = glob.glob(f'{TARGET_DIR}/force/*.png')
energy_pic.sort()
force_pic.sort()

IMG_DISPLAY_HEIGHT = Inches(2.5)  # 画像サイズ
PIC_PER_PAGE = 2
OUTPUT_FILE_NAME = f"/Users/y1u0d2/desktop/Lab/result/nnp-train/20211130/r2.pptx"  # 出力ファイル名

slides=len(energy_pic)
# if len(energy_pic) % 2 == 0:
#     slides = len(energy_pic)
# else:
#     slides = (len(energy_pic)+1)/2

prs = Presentation()
SLIDE_WIDTH = prs.slide_width
SLIDE_HEIGHT = prs.slide_height

for i in range(int(slides)):
    # slide_layouts:6 blank
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    energy_pic_1 = energy_pic[i]
    force_pic_1 = force_pic[i]

    im = Image.open(energy_pic_1)
    im_width, im_height = im.size
    aspect_ratio = im_width / im_height

    img_display_height = IMG_DISPLAY_HEIGHT
    img_display_width = aspect_ratio * img_display_height

    # top
    left = (SLIDE_WIDTH - img_display_width)/2
    top = (SLIDE_HEIGHT / 2 - img_display_height)
    slide.shapes.add_picture(energy_pic_1, left, top, height=IMG_DISPLAY_HEIGHT)

    # bottom
    left = (SLIDE_WIDTH - img_display_width)/2
    top = (SLIDE_HEIGHT / 2)
    slide.shapes.add_picture(force_pic_1, left, top, height=IMG_DISPLAY_HEIGHT)

prs.save(OUTPUT_FILE_NAME)
# TODO convert pdf