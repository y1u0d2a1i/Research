from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import glob

from utils.constants.constants import Constants


structures = Constants.structures()
TARGET_DIR = '/Users/y1u0d2/desktop/Lab/result/lammps/structure-optimization/10'
all_pic = []
for structure in structures:
    dirs = glob.glob(f'{TARGET_DIR}/{structure}/*')
    for _dir in dirs:
        pic_min = glob.glob(f'{_dir}/*_min.png')[0]
        pic_rdf = glob.glob(f'{_dir}/{structure}_*-*.png')
        all_pic.append(pic_min)
        for pic in pic_rdf:
            all_pic.append(pic)

IMG_DISPLAY_HEIGHT = Inches(3.5)  # 画像サイズ
PIC_PER_PAGE = 4
OUTPUT_FILE_NAME = '/Users/y1u0d2/desktop/Lab/result/lammps/structure-optimization/10/rdf_comparison.pptx'  # 出力ファイル名

slides=len(all_pic)/4

prs = Presentation()
SLIDE_WIDTH = prs.slide_width
SLIDE_HEIGHT = prs.slide_height

for i in range(int(slides)):
    # slide_layouts:6 blank
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic1 = all_pic[i*4]
    pic2 = all_pic[i*4+1]
    pic3 = all_pic[i*4+2]
    pic4 = all_pic[i*4+3]

    im = Image.open(pic2)
    im_width, im_height = im.size
    aspect_ratio = im_width / im_height

    img_display_height = IMG_DISPLAY_HEIGHT
    img_display_width = aspect_ratio * img_display_height

    # left top
    left = (SLIDE_WIDTH / 2 - img_display_width)
    top = (SLIDE_HEIGHT / 2 - img_display_height)
    slide.shapes.add_picture(pic1, left, top, height=IMG_DISPLAY_HEIGHT)
    # left bottom
    if not i*2+1 >= len(all_pic):
        left = (SLIDE_WIDTH / 2 - img_display_width)
        top = (SLIDE_HEIGHT / 2)
        slide.shapes.add_picture(pic2, left, top, height=IMG_DISPLAY_HEIGHT)
    # right top
    left = (SLIDE_WIDTH / 2)
    top = (SLIDE_HEIGHT / 2 - img_display_height)
    slide.shapes.add_picture(pic3, left, top, height=IMG_DISPLAY_HEIGHT)
    # right bottom
    if not i*2+1 >= len(all_pic):
        left = (SLIDE_WIDTH / 2)
        top = (SLIDE_HEIGHT / 2)
        slide.shapes.add_picture(pic4, left, top, height=IMG_DISPLAY_HEIGHT)

prs.save(OUTPUT_FILE_NAME)